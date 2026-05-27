#!/usr/bin/env python3
"""
PPT Translator — On-Prem Presentation Translation
===================================================
Upload a .pptx → LLM translates all text slide-by-slide → download translated .pptx.
Preserves all formatting, images, layouts, and styles.

Runs 100% on-prem on the HP ZGX Nano via vLLM.
"""

import asyncio
import json
import logging
import os
import shutil
import time
import uuid
from copy import deepcopy
from pathlib import Path
from datetime import datetime

import httpx
from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Request
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
from pptx.util import Pt, Emu
from pydantic import BaseModel

# ── Config ────────────────────────────────────────────────────────────────────

BASE_DIR = Path(__file__).parent
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "output"
TEMPLATES_DIR = BASE_DIR / "templates"
STATIC_DIR = BASE_DIR / "static"

UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

VLLM_URL = os.getenv("VLLM_URL", "http://localhost:8091")
VLLM_MODEL = os.getenv("VLLM_MODEL", "Qwen/Qwen3.6-35B-A3B")

# ── Logging ───────────────────────────────────────────────────────────────────

logging.basicConfig(
    level=logging.INFO,
    format="[%(asctime)s] [%(name)s] [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
log = logging.getLogger("PPT-Translator")

# ── Language Map ──────────────────────────────────────────────────────────────

LANGUAGES = {
    "ar": "Arabic",
    "zh": "Chinese (Simplified)",
    "zh-tw": "Chinese (Traditional)",
    "cs": "Czech",
    "da": "Danish",
    "nl": "Dutch",
    "fi": "Finnish",
    "fr": "French",
    "de": "German",
    "hi": "Hindi",
    "hu": "Hungarian",
    "id": "Indonesian",
    "it": "Italian",
    "ja": "Japanese",
    "ko": "Korean",
    "ms": "Malay",
    "nb": "Norwegian Bokmål",
    "pl": "Polish",
    "pt": "Portuguese (Brazil)",
    "pt-pt": "Portuguese (Portugal)",
    "ro": "Romanian",
    "ru": "Russian",
    "es": "Spanish",
    "sv": "Swedish",
    "th": "Thai",
    "tr": "Turkish",
    "uk": "Ukrainian",
    "vi": "Vietnamese",
}

# ── App ───────────────────────────────────────────────────────────────────────

app = FastAPI(title="PPT Translator")


# ── PPTX Text Extraction ─────────────────────────────────────────────────────

def _iter_shapes(shapes):
    """
    Recursively yield all shapes, descending into group shapes.
    Handles arbitrarily nested groups.
    """
    for shape in shapes:
        yield shape
        if shape.shape_type is not None and shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                yield from _iter_shapes(shape.shapes)
            except Exception:
                pass


def extract_slide_texts(prs: Presentation) -> list[dict]:
    """
    Extract all text from each slide, preserving shape-level grouping.
    Returns a list of dicts: [{"slide_index": 0, "texts": ["text1", "text2", ...]}, ...]
    Recursively descends into group shapes.
    Only includes non-empty text runs.
    """
    slides_data = []
    for idx, slide in enumerate(prs.slides):
        texts = []
        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip()
                if full_text:
                    texts.append(full_text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            texts.append(cell_text)
        if texts:
            slides_data.append({"slide_index": idx, "texts": texts})
    return slides_data


def apply_translated_text(prs: Presentation, slide_index: int, translations: dict[str, str]):
    """
    Walk through shapes on the given slide and replace original text with translations.
    Preserves all formatting (font, size, color, bold, italic, etc.) by replacing
    text at the run level while keeping run properties intact.
    
    translations is a dict mapping original_text -> translated_text.
    """
    slide = prs.slides[slide_index]
    for shape in slide.shapes:
        if shape.has_text_frame:
            _replace_textframe(shape.text_frame, translations)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _replace_textframe(cell.text_frame, translations)


def _replace_textframe(text_frame, translations: dict[str, str]):
    """
    Replace text in a text frame while preserving formatting.
    Strategy: for each paragraph, concatenate all runs to get the full paragraph text.
    If that text matches a translation key, replace it by putting all translated text
    into the first run and clearing subsequent runs (preserving the first run's formatting).
    """
    for paragraph in text_frame.paragraphs:
        # Get full paragraph text by joining runs
        full_para_text = "".join(run.text for run in paragraph.runs).strip()
        if not full_para_text:
            continue

        # Check if this paragraph text (or the shape-level text) has a translation
        translated = None
        for orig, trans in translations.items():
            if full_para_text == orig:
                translated = trans
                break

        if translated is None:
            # Try matching as a substring of a larger text block
            # This handles multi-paragraph shapes where we translated the whole block
            continue

        # Apply translation preserving formatting of first run
        runs = paragraph.runs
        if not runs:
            continue

        # Put all translated text into the first run
        runs[0].text = translated
        # Clear remaining runs
        for run in runs[1:]:
            run.text = ""


def apply_translations_whole_shape(prs: Presentation, slide_index: int, translations: dict[str, str]):
    """
    Alternative approach: match at the whole-shape level.
    For each shape (recursing into groups), get the full text, look up translation, then redistribute.
    """
    slide = prs.slides[slide_index]
    for shape in _iter_shapes(slide.shapes):
        if shape.has_text_frame:
            shape_text = shape.text_frame.text.strip()
            if shape_text in translations:
                _apply_to_shape_textframe(shape.text_frame, shape_text, translations[shape_text])
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text in translations:
                        _apply_to_shape_textframe(cell.text_frame, cell_text, translations[cell_text])


def _apply_to_shape_textframe(text_frame, original: str, translated: str):
    """
    Replace entire text frame content with translated text.
    Preserves formatting of the first run of the first paragraph.
    For multi-paragraph content, splits translated text by newlines and maps
    to existing paragraphs where possible.
    """
    orig_paragraphs = original.split("\n")
    trans_paragraphs = translated.split("\n")

    paragraphs = list(text_frame.paragraphs)

    # If same number of paragraphs (or close), do 1:1 mapping
    for i, paragraph in enumerate(paragraphs):
        para_text = "".join(run.text for run in paragraph.runs).strip()
        if not para_text:
            continue

        # Find which translated paragraph maps here
        # Simple approach: map by index within non-empty paragraphs
        if i < len(trans_paragraphs):
            new_text = trans_paragraphs[i]
        else:
            new_text = ""

        runs = paragraph.runs
        if runs:
            runs[0].text = new_text
            for run in runs[1:]:
                run.text = ""

    # If we have more translated paragraphs than original, append remaining
    # to the last paragraph's first run
    if len(trans_paragraphs) > len(paragraphs) and paragraphs:
        last_para = paragraphs[-1]
        if last_para.runs:
            extra = "\n".join(trans_paragraphs[len(paragraphs):])
            last_para.runs[0].text += "\n" + extra


# ── vLLM Translation ─────────────────────────────────────────────────────────

async def translate_texts(texts: list[str], target_lang: str, source_lang: str = "English") -> list[str]:
    """
    Send a batch of text strings to vLLM for translation.
    Uses a single prompt with numbered items for efficiency.
    Returns translated strings in the same order.
    """
    if not texts:
        return []

    target_lang_name = LANGUAGES.get(target_lang, target_lang)

    # Build numbered list
    numbered = "\n".join(f"[{i+1}] {t}" for i, t in enumerate(texts))

    prompt = f"""You are a professional translator specializing in corporate training materials and business documents.

Translate the following numbered text items from {source_lang} to {target_lang_name}.

CRITICAL RULES:
- Return ONLY a JSON array of strings, one per numbered item, in the same order
- Preserve any formatting markers, numbers, or special characters
- Keep proper nouns, brand names, product names, and technical terms (e.g. "HP", "ZGX Nano", "Windows", "PowerPoint") untranslated
- Use formal/professional register appropriate for corporate training materials
- Do NOT add explanations, notes, or commentary
- Do NOT include the original text
- Do NOT include numbering in the output

Text items to translate:

{numbered}

Respond with ONLY a JSON array like: ["translated item 1", "translated item 2", ...]"""

    payload = {
        "model": VLLM_MODEL,
        "messages": [
            {"role": "user", "content": prompt}
        ],
        "max_tokens": 8192,
        "temperature": 0.3,
    }

    try:
        async with httpx.AsyncClient(timeout=120) as client:
            resp = await client.post(f"{VLLM_URL}/v1/chat/completions", json=payload)
            resp.raise_for_status()

            data = resp.json()
            content = data["choices"][0]["message"]["content"]

            # Strip thinking tags if present (Qwen3 /think mode)
            if "<think>" in content:
                # Remove everything between <think> and </think>
                import re
                content = re.sub(r"<think>.*?</think>", "", content, flags=re.DOTALL).strip()

            # Parse JSON array from response
            # Try to find JSON array in the response
            content = content.strip()
            if content.startswith("```"):
                # Remove code fences
                content = content.split("\n", 1)[1] if "\n" in content else content[3:]
                if content.endswith("```"):
                    content = content[:-3]
                content = content.strip()
                if content.startswith("json"):
                    content = content[4:].strip()

            translations = json.loads(content)

            if not isinstance(translations, list):
                raise ValueError(f"Expected JSON array, got {type(translations)}")

            if len(translations) != len(texts):
                log.warning(
                    f"Translation count mismatch: expected {len(texts)}, got {len(translations)}. "
                    f"Padding/truncating."
                )
                # Pad with originals if too few
                while len(translations) < len(texts):
                    translations.append(texts[len(translations)])
                # Truncate if too many
                translations = translations[:len(texts)]

            return translations

    except json.JSONDecodeError as e:
        log.error(f"Failed to parse translation JSON: {e}\nContent: {content[:500]}")
        raise HTTPException(status_code=500, detail=f"Translation response was not valid JSON: {e}")
    except httpx.HTTPStatusError as e:
        log.error(f"vLLM returned error: {e.response.status_code} - {e.response.text[:300]}")
        raise HTTPException(status_code=502, detail=f"Inference engine error: {e.response.status_code}")
    except httpx.ConnectError:
        log.error(f"Cannot connect to vLLM at {VLLM_URL}")
        raise HTTPException(status_code=503, detail="Inference engine not reachable. Is vLLM running?")
    except Exception as e:
        log.error(f"Translation error: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"Translation failed: {e}")


# ── Job Tracking ──────────────────────────────────────────────────────────────

jobs: dict[str, dict] = {}


# ── API Endpoints ─────────────────────────────────────────────────────────────

@app.post("/api/translate")
async def translate_pptx(
    file: UploadFile = File(...),
    target_lang: str = Form(...),
    source_lang: str = Form(default="en"),
):
    """Upload a PPTX and translate all text content."""
    job_id = f"job_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:6]}"
    log.info(f"[{job_id}] New translation job: {file.filename} → {target_lang}")

    # Validate file type
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported")

    # Save upload
    upload_path = UPLOAD_DIR / f"{job_id}_original.pptx"
    content = await file.read()
    upload_path.write_bytes(content)
    log.info(f"[{job_id}] Saved upload: {upload_path.name} ({len(content):,} bytes)")

    # Track job
    jobs[job_id] = {
        "status": "processing",
        "filename": file.filename,
        "target_lang": target_lang,
        "source_lang": source_lang,
        "total_slides": 0,
        "translated_slides": 0,
        "total_texts": 0,
        "translated_texts": 0,
        "started_at": time.time(),
        "error": None,
    }

    # Process in background
    asyncio.create_task(_process_translation(job_id, upload_path, target_lang, source_lang, file.filename))

    return {"job_id": job_id, "status": "processing"}


async def _process_translation(job_id: str, upload_path: Path, target_lang: str, source_lang: str, original_filename: str):
    """Background task: extract → translate → rebuild PPTX."""
    try:
        prs = Presentation(str(upload_path))
        total_slides = len(prs.slides)
        jobs[job_id]["total_slides"] = total_slides
        log.info(f"[{job_id}] Loaded presentation: {total_slides} slides")

        # Extract text from all slides
        slides_data = extract_slide_texts(prs)
        total_texts = sum(len(sd["texts"]) for sd in slides_data)
        jobs[job_id]["total_texts"] = total_texts
        log.info(f"[{job_id}] Extracted {total_texts} text blocks from {len(slides_data)} slides with content")

        source_lang_name = LANGUAGES.get(source_lang, "English")
        if source_lang == "en":
            source_lang_name = "English"

        # Translate slide by slide
        for sd in slides_data:
            slide_idx = sd["slide_index"]
            texts = sd["texts"]

            log.info(f"[{job_id}] Translating slide {slide_idx + 1}: {len(texts)} text blocks")

            # Translate this slide's texts
            translated = await translate_texts(texts, target_lang, source_lang_name)

            # Build translation map
            translation_map = dict(zip(texts, translated))

            # Apply to presentation
            apply_translations_whole_shape(prs, slide_idx, translation_map)

            jobs[job_id]["translated_slides"] += 1
            jobs[job_id]["translated_texts"] += len(texts)
            log.info(f"[{job_id}] Slide {slide_idx + 1} translated")

        # Save translated PPTX
        output_path = OUTPUT_DIR / f"{job_id}.pptx"
        prs.save(str(output_path))

        file_size = output_path.stat().st_size
        elapsed = time.time() - jobs[job_id]["started_at"]

        jobs[job_id]["status"] = "complete"
        jobs[job_id]["elapsed"] = round(elapsed, 2)
        jobs[job_id]["output_size"] = file_size
        jobs[job_id]["download_url"] = f"/api/download/{job_id}"
        jobs[job_id]["output_filename"] = _make_output_filename(original_filename, target_lang)

        log.info(f"[{job_id}] Translation complete: {elapsed:.1f}s, {file_size:,} bytes")

    except Exception as e:
        log.error(f"[{job_id}] Translation failed: {type(e).__name__}: {e}")
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)


def _make_output_filename(original: str, target_lang: str) -> str:
    """Generate output filename like 'training_deck_fr.pptx'."""
    stem = Path(original).stem
    lang_name = LANGUAGES.get(target_lang, target_lang)
    short_lang = target_lang.replace("-", "")
    return f"{stem}_{short_lang}.pptx"


@app.get("/api/status/{job_id}")
async def get_job_status(job_id: str):
    """Poll translation job status."""
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    return jobs[job_id]


@app.get("/api/download/{job_id}")
async def download_translated(job_id: str):
    """Download the translated PPTX."""
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")

    job = jobs[job_id]
    if job["status"] != "complete":
        raise HTTPException(status_code=400, detail=f"Job is not complete (status: {job['status']})")

    output_path = OUTPUT_DIR / f"{job_id}.pptx"
    if not output_path.exists():
        raise HTTPException(status_code=404, detail="Output file not found")

    return FileResponse(
        path=str(output_path),
        filename=job.get("output_filename", f"{job_id}.pptx"),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@app.get("/api/languages")
async def get_languages():
    """Return available target languages."""
    return {"languages": LANGUAGES}


@app.get("/api/health")
async def health():
    """Health check including vLLM status."""
    vllm_ok = False
    vllm_model = None
    try:
        async with httpx.AsyncClient(timeout=5) as client:
            r = await client.get(f"{VLLM_URL}/health")
            vllm_ok = r.status_code == 200
            if vllm_ok:
                models_resp = await client.get(f"{VLLM_URL}/v1/models")
                if models_resp.status_code == 200:
                    model_data = models_resp.json()
                    loaded = [m["id"] for m in model_data.get("data", [])]
                    vllm_model = loaded[0] if loaded else None
    except Exception:
        pass

    return {
        "status": "ok" if vllm_ok else "degraded",
        "vllm": "healthy" if vllm_ok else "unreachable",
        "vllm_url": VLLM_URL,
        "model": vllm_model or VLLM_MODEL,
        "active_jobs": sum(1 for j in jobs.values() if j["status"] == "processing"),
    }


# ── Frontend ──────────────────────────────────────────────────────────────────

@app.get("/")
async def serve_frontend():
    index_path = TEMPLATES_DIR / "index.html"
    if index_path.exists():
        return HTMLResponse(content=index_path.read_text())
    return HTMLResponse(content="<h1>PPT Translator</h1><p>Frontend not found.</p>")


# Mount static files
if STATIC_DIR.exists():
    app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")


# ── Entry Point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8092, log_level="info")